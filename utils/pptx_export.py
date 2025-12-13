"""
PowerPoint export utility for QC Presentation Tool.

Creates professional presentations with seismic images and spectral plots.
"""
import io
import numpy as np
from typing import Optional
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from PIL import Image
import matplotlib.pyplot as plt
from matplotlib.figure import Figure


class PPTXExporter:
    """
    Export QC Presentation session to PowerPoint.

    Creates slides with:
    - Large seismic/spectral image (centered)
    - Title (band + view type)
    - Subtitle (gather info, gain, parameters)
    """

    def __init__(self, output_path: str):
        """
        Initialize exporter.

        Args:
            output_path: Path for output .pptx file
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is not installed. "
                "Install with: pip install python-pptx"
            )

        self.output_path = output_path
        self.prs = Presentation()

        # Set widescreen 16:9 format
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

    def add_seismic_slide(self,
                          image_array: np.ndarray,
                          title: str,
                          subtitle: str = "",
                          colormap: str = 'seismic',
                          sample_rate: float = 4.0):
        """
        Add slide with seismic image.

        Args:
            image_array: 2D numpy array (n_samples, n_traces)
            title: Slide title
            subtitle: Slide subtitle
            colormap: Matplotlib colormap name
            sample_rate: Sample rate in ms (for axis labels)
        """
        # Create blank slide
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # Create seismic image with matplotlib
        fig = self._create_seismic_figure(image_array, colormap, sample_rate)
        img_buffer = self._figure_to_buffer(fig)
        plt.close(fig)

        # Add image (large, centered)
        left = Inches(0.5)
        top = Inches(1.0)
        width = Inches(12.333)
        height = Inches(5.5)
        slide.shapes.add_picture(img_buffer, left, top, width, height)

        # Add title
        self._add_title(slide, title)

        # Add subtitle
        if subtitle:
            self._add_subtitle(slide, subtitle)

    def add_spectrum_slide(self,
                           figure: Figure,
                           title: str,
                           subtitle: str = ""):
        """
        Add slide with matplotlib spectrum figure.

        Args:
            figure: Matplotlib Figure object
            title: Slide title
            subtitle: Slide subtitle
        """
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # Save figure to buffer
        img_buffer = self._figure_to_buffer(figure)

        # Add image
        left = Inches(1.0)
        top = Inches(1.0)
        width = Inches(11.333)
        height = Inches(5.5)
        slide.shapes.add_picture(img_buffer, left, top, width, height)

        # Add title
        self._add_title(slide, title)

        # Add subtitle
        if subtitle:
            self._add_subtitle(slide, subtitle)

    def add_title_slide(self,
                        title: str,
                        subtitle: str = ""):
        """
        Add a title slide.

        Args:
            title: Main title
            subtitle: Subtitle
        """
        title_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(title_layout)

        # Centered title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.333), Inches(1.0)
            )
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = subtitle
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = RGBColor(100, 100, 100)
            sub_para.alignment = PP_ALIGN.CENTER

    def save(self):
        """Save presentation to disk."""
        self.prs.save(self.output_path)

    def _create_seismic_figure(self,
                               data: np.ndarray,
                               colormap: str,
                               sample_rate: float) -> Figure:
        """Create matplotlib figure for seismic data."""
        fig, ax = plt.subplots(figsize=(12, 8), facecolor='white')

        # Compute time axis
        n_samples, n_traces = data.shape
        time_max = n_samples * sample_rate

        # Plot seismic with proper orientation
        extent = [0, n_traces - 1, time_max, 0]  # Time increases downward

        # Normalize for display
        clip_val = np.percentile(np.abs(data), 99)
        if clip_val > 0:
            vmin, vmax = -clip_val, clip_val
        else:
            vmin, vmax = -1, 1

        im = ax.imshow(
            data,
            aspect='auto',
            cmap=colormap,
            extent=extent,
            vmin=vmin,
            vmax=vmax,
            interpolation='bilinear'
        )

        ax.set_xlabel('Trace Number', fontsize=12)
        ax.set_ylabel('Time (ms)', fontsize=12)

        # Add colorbar
        cbar = fig.colorbar(im, ax=ax, shrink=0.8)
        cbar.set_label('Amplitude', fontsize=10)

        fig.tight_layout()
        return fig

    def _figure_to_buffer(self, figure: Figure) -> io.BytesIO:
        """Convert matplotlib figure to PNG buffer."""
        buf = io.BytesIO()
        figure.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                       facecolor='white', edgecolor='none')
        buf.seek(0)
        return buf

    def _array_to_image(self, arr: np.ndarray, colormap: str) -> io.BytesIO:
        """Convert numpy array directly to PNG image buffer."""
        # Normalize to 0-1
        arr_norm = arr - arr.min()
        max_val = arr_norm.max()
        if max_val > 0:
            arr_norm = arr_norm / max_val

        # Apply colormap
        cmap = plt.get_cmap(colormap)
        colored = cmap(arr_norm)
        colored = (colored[:, :, :3] * 255).astype(np.uint8)

        # Convert to PIL and save to buffer
        img = Image.fromarray(colored)
        buf = io.BytesIO()
        img.save(buf, format='PNG')
        buf.seek(0)
        return buf

    def _add_title(self, slide, title: str):
        """Add title to slide."""
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER

    def _add_subtitle(self, slide, subtitle: str):
        """Add subtitle to slide."""
        sub_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.7), Inches(12.333), Inches(0.4)
        )
        sub_frame = sub_box.text_frame
        sub_para = sub_frame.paragraphs[0]
        sub_para.text = subtitle
        sub_para.font.size = Pt(14)
        sub_para.font.color.rgb = RGBColor(100, 100, 100)
        sub_para.alignment = PP_ALIGN.CENTER


def is_pptx_available() -> bool:
    """Check if python-pptx is available."""
    return PPTX_AVAILABLE
